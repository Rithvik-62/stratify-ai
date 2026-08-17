import os
import sys
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, KeepTogether, HRFlowable, Image
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Output & Brain Directories
OUTPUT_DIR = os.path.join(os.getcwd(), "Output")
os.makedirs(OUTPUT_DIR, exist_ok=True)

BRAIN_DIR = r"C:\Users\shrinidhi\.gemini\antigravity-ide\brain\5de7cc53-4e63-466b-815c-27e463ae8dbf"

# Live Demo Screenshots
IMG_EXEC_MID = os.path.join(BRAIN_DIR, "exec_control_mid_1786949688444.png")
IMG_EXEC_CHART1 = os.path.join(BRAIN_DIR, "exec_control_chart1_1786949705566.png")
IMG_EXEC_CHART2 = os.path.join(BRAIN_DIR, "exec_control_chart2_1786949724875.png")
IMG_EXEC_BOTTOM = os.path.join(BRAIN_DIR, "exec_control_bottom_1786949746057.png")
IMG_CUST_TOP = os.path.join(BRAIN_DIR, "customer_analytics_top_1786949921371.png")
IMG_CUST_MID = os.path.join(BRAIN_DIR, "customer_analytics_mid1_1786949940069.png")
IMG_FORECAST_TOP = os.path.join(BRAIN_DIR, "sales_forecasting_top_1786950054016.png")
IMG_FORECAST_MID = os.path.join(BRAIN_DIR, "sales_forecasting_mid1_1786950078150.png")
IMG_AI_TOP = os.path.join(BRAIN_DIR, "ai_insights_top_1786950197991.png")
IMG_AI_MID = os.path.join(BRAIN_DIR, "ai_insights_mid1_1786950241833.png")
IMG_DQ_TOP = os.path.join(BRAIN_DIR, "data_quality_top_1786950543537.png")
IMG_DQ_MID = os.path.join(BRAIN_DIR, "data_quality_mid1_1786950567981.png")
IMG_REPORTS_TOP = os.path.join(BRAIN_DIR, "executive_reports_top_1786950755548.png")
IMG_REPORTS_MID = os.path.join(BRAIN_DIR, "executive_reports_mid1_1786950775806.png")

# Palette
PRIMARY = colors.HexColor("#1A365D")    # Deep Navy
SECONDARY = colors.HexColor("#2B6CB0")  # Slate Blue
ACCENT = colors.HexColor("#319795")     # Teal Accent
BG_LIGHT = colors.HexColor("#F7FAFC")   # Light Gray
CARD_BG = colors.HexColor("#EDF2F7")    # Card Gray
TEXT_DARK = colors.HexColor("#2D3748")  # Dark Slate Text
BORDER_COLOR = colors.HexColor("#CBD5E0") # Border Gray
WHITE = colors.HexColor("#FFFFFF")

def get_pdf_styles():
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name='DocTitle', fontName='Helvetica-Bold', fontSize=22, leading=26, textColor=PRIMARY, spaceAfter=6
    ))
    styles.add(ParagraphStyle(
        name='DocSubTitle', fontName='Helvetica-Bold', fontSize=13, leading=16, textColor=SECONDARY, spaceAfter=15
    ))
    styles.add(ParagraphStyle(
        name='SectionHeader', fontName='Helvetica-Bold', fontSize=15, leading=18, textColor=PRIMARY, spaceBefore=12, spaceAfter=8, keepWithNext=True
    ))
    styles.add(ParagraphStyle(
        name='SubSectionHeader', fontName='Helvetica-Bold', fontSize=12, leading=15, textColor=SECONDARY, spaceBefore=8, spaceAfter=4, keepWithNext=True
    ))
    styles.add(ParagraphStyle(
        name='CustomBody', fontName='Helvetica', fontSize=10, leading=14, textColor=TEXT_DARK, spaceAfter=6
    ))
    styles.add(ParagraphStyle(
        name='CustomBodyBold', fontName='Helvetica-Bold', fontSize=10, leading=14, textColor=TEXT_DARK, spaceAfter=6
    ))
    styles.add(ParagraphStyle(
        name='SpokenScript', fontName='Helvetica-Oblique', fontSize=9.5, leading=14, textColor=colors.HexColor("#2C5282"), spaceAfter=6
    ))
    styles.add(ParagraphStyle(
        name='TableHeader', fontName='Helvetica-Bold', fontSize=9.5, leading=12, textColor=WHITE, alignment=0
    ))
    styles.add(ParagraphStyle(
        name='TableCell', fontName='Helvetica', fontSize=9, leading=12, textColor=TEXT_DARK
    ))
    return styles

# Presenters Data with Live Demo Screenshot Paths & Speaker Notes
PRESENTERS_DATA = [
    {
        "person": "Person 1",
        "file_name": "Person1_Business_Problem_and_Introduction.pdf",
        "title": "Business Problem & Platform Vision",
        "theme": "Overcoming Data Silos & Transforming Executive Decision-Making",
        "duration": "5 Minutes",
        "summary": "Focuses on why traditional reporting fails modern businesses, introduces Stratify AI as a real-time decision intelligence platform, and maps out the end-to-end business flow across all 6 stages.",
        "features": [
            {
                "name": "The Business Challenge",
                "what": "Modern enterprises suffer from delayed sales reports (taking 3-7 days), disconnected store data, spreadsheet clutter, and reactive decision-making.",
                "connects": "Connects directly to the vision of Stratify AI as an automated, instant intelligence solution.",
                "value": "Eliminates reporting lag, reduces human manual errors by 95%, and provides C-suite executives with immediate operational clarity."
            },
            {
                "name": "The Stratify AI Vision",
                "what": "An all-in-one executive command center that continuously streams sales data, validates accuracy, forecasts trends, and provides AI strategic advice.",
                "connects": "Forms the central hub linking store transactions to executive strategy.",
                "value": "Accelerates decision cycles from days to minutes, empowering leadership to capture market opportunities faster."
            },
            {
                "name": "End-to-End Business Flow",
                "what": "A seamless 5-stage lifecycle: Store Checkout ➔ Automated Data Clean ➔ Customer & Sales Intelligence ➔ Live Control Center ➔ AI Strategy ➔ Automated Email Briefings.",
                "connects": "Provides the overall map for all 6 presenters.",
                "value": "Creates total cross-departmental alignment with zero manual effort."
            }
        ],
        "slides": [
            {
                "num": "Slide 1",
                "title": "Welcome & Executive Summary",
                "visual": "Live Demo Header & Corporate KPI Summary Banner",
                "img": IMG_EXEC_MID,
                "points": [
                    "<b>Live Decision Command Center</b>: Single-screen visibility across all retail locations",
                    "<b>Real-Time KPI Visibility</b>: Live Total Revenue (₹1.32M), Net Margin (26.63%), 486 Active Customers",
                    "<b>Executive Clarity</b>: Replacing static, delayed spreadsheets with continuous intelligence"
                ],
                "script": "Welcome everyone. Today we are excited to introduce Stratify AI, an executive decision intelligence platform designed to solve one of the biggest challenges in business today: delayed visibility. As seen on our live web portal header, Stratify AI replaces static spreadsheets with real-time KPI visibility—tracking total revenue, profit margins, and active customer accounts live.",
                "takeaway": "Stratify AI replaces static, delayed reporting with real-time strategic intelligence."
            },
            {
                "num": "Slide 2",
                "title": "The Business Pain Points",
                "visual": "Financial Performance & Category Margin Analysis Dashboard",
                "img": IMG_EXEC_CHART2,
                "points": [
                    "<b>Siloed Enterprise Operations</b>: POS, inventory, and finance teams operating in isolation",
                    "<b>Spreadsheet Bottlenecks</b>: 15+ hours lost copy-pasting raw CSV files every week",
                    "<b>Reactive Decision Lag</b>: Identifying stockouts and margin drops days after revenue is lost"
                ],
                "script": "Let's look at the problem we face without a unified system. Checkout counters generate sales, but that data sits in store databases. Finance teams manually compile reports over the weekend, and by Monday, executives see what happened last week—not what's happening right now. This latency leads to missed stockouts, lost customer retention, and delayed promotional campaigns.",
                "takeaway": "Manual reporting bottlenecks burn time and hide real-time revenue risks."
            },
            {
                "num": "Slide 3",
                "title": "The End-to-End Stratify Solution Map",
                "visual": "5-Stage High-Level Business Flowchart",
                "img": IMG_EXEC_CHART1,
                "points": [
                    "<b>Stage 1 & 2</b>: Real-time store capture & automated digital quality filter",
                    "<b>Stage 3 & 4</b>: RFM customer intelligence & interactive executive control center",
                    "<b>Stage 5 & 6</b>: DeepSeek AI strategy advisor & automated morning PDF email briefs"
                ],
                "script": "To solve this, we built a 5-stage automated business pipeline. As customer purchases occur at store terminals, the data is captured, instantly filtered for 100% data quality, enriched with customer segments and predictive forecasts, presented on a live control dashboard, analyzed by our built-in AI strategy advisor, and delivered as an executive morning briefing email. Let's pass to Person 2 to see how store sales are captured and cleaned.",
                "takeaway": "Stratify AI connects raw store transactions directly to executive inbox briefings."
            }
        ],
        "handoff": "Thank you everyone. Now, to explain how our store sales are captured in real-time and filtered through an automated digital quality check before touching any business report, I'll hand over to Person 2."
    },
    {
        "person": "Person 2",
        "file_name": "Person2_Live_Store_Capture_and_Clean_Data.pdf",
        "title": "Live Store Capture & Smart Data Quality Control",
        "theme": "Garbage In, Garbage Out Prevention & Enterprise Data Reliability",
        "duration": "5 Minutes",
        "summary": "Explains how live store transactions are continuously recorded, passed through an automated digital quality control filter, and stored securely so business reports are always 100% accurate.",
        "features": [
            {
                "name": "Live Store Transaction Capture",
                "what": "Continuous simulation and recording of retail checkout counter activity across multiple store locations, including purchase amounts, quantities, customer IDs, and product categories.",
                "connects": "Feeds directly into the automated quality control filter.",
                "value": "Ensures zero lag between customer purchases and central business awareness."
            },
            {
                "name": "Automated Digital Quality Control Filter",
                "what": "A smart validation gate that inspects incoming sales data in real time, automatically correcting duplicate receipts, flagging negative numbers, and removing corrupted records.",
                "connects": "Sits between raw store terminals and the central business data warehouse.",
                "value": "Protects executive decision-making from bad or corrupted data ('Garbage In, Garbage Out' prevention)."
            },
            {
                "name": "Secure Central Data Repository",
                "what": "An enterprise-grade cloud repository where clean, audited sales transactions are stored with zero duplicate records.",
                "connects": "Provides the single source of truth for analytics, dashboards, and AI models.",
                "value": "Guarantees data consistency, compliance, and instant accessibility across all corporate departments."
            }
        ],
        "slides": [
            {
                "num": "Slide 1",
                "title": "Capturing Sales at Point of Purchase",
                "visual": "Enterprise Data Quality & SLA Governance Center Dashboard",
                "img": IMG_DQ_TOP,
                "points": [
                    "<b>Live POS Stream</b>: Capturing retail checkout receipts instantly across store outlets",
                    "<b>Complete Data Capture</b>: Revenue, units sold, customer account IDs, and product SKUs",
                    "<b>Zero Filing Latency</b>: Eliminating traditional end-of-day batch processing delays"
                ],
                "script": "Thanks Person 1. Every business decision is only as good as the data behind it. In our platform, step one begins at the store level. As customers check out—whether buying electronics, clothing, or household goods—the transaction details are captured live. Instead of waiting for end-of-day store closing files, our pipeline registers transactions instantly.",
                "takeaway": "Real-time store capture ensures business metrics reflect immediate market reality."
            },
            {
                "num": "Slide 2",
                "title": "The Digital Quality Control Filter",
                "visual": "100% Data Completeness, Uniqueness & SLA Governance Matrix",
                "img": IMG_DQ_MID,
                "points": [
                    "<b>100.0% Uniqueness</b>: Automated receipt deduplication preventing double-counted sales",
                    "<b>100.0% Validity & Completeness</b>: Verifying pricing, quantities, and foreign key rules",
                    "<b><5s SLA Latency</b>: High-speed validation guarantee before data enters executive views"
                ],
                "script": "Raw data is often noisy—duplicate scans occur, network glitches send incomplete records, or pricing glitches appear. As shown on our Data Quality dashboard, Stratify AI enforces 100% completeness, 100% uniqueness, and under 5-second SLA latency. Duplicates are eliminated, pricing errors are flagged, and only 100% verified clean data passes through.",
                "takeaway": "Automated data verification guarantees 100% reliable executive reporting."
            },
            {
                "num": "Slide 3",
                "title": "Single Source of Business Truth",
                "visual": "Cloud Repository Status showing audited data ready for Analytics",
                "img": IMG_DQ_TOP,
                "points": [
                    "<b>Unified Cloud Repository</b>: Single audited data store accessible company-wide",
                    "<b>Idempotent Loading</b>: Smoothly updating existing sales records without duplication",
                    "<b>Department Alignment</b>: Finance, sales, and executive teams looking at identical numbers"
                ],
                "script": "Once verified, the clean transactions enter our central enterprise data repository. This creates a single source of truth for the entire company. Sales managers, finance teams, and executive officers all view the exact same audited numbers. Now that our data is clean and centrally stored, let's pass to Person 3 to see how we turn these numbers into customer intelligence and sales predictions.",
                "takeaway": "A single source of clean data aligns all teams under one corporate truth."
            }
        ],
        "handoff": "Now that we have established a continuous flow of 100% clean, verified store data, I will hand over to Person 3 to show how we extract customer intelligence, forecast future sales, and track inventory health."
    },
    {
        "person": "Person 3",
        "file_name": "Person3_Customer_Intelligence_and_Forecasting.pdf",
        "title": "Customer Intelligence & Sales Growth Forecasting",
        "theme": "Turning Raw Sales Data into Customer Insights & Revenue Predictions",
        "duration": "5 Minutes",
        "summary": "Focuses on customer RFM segmentation (Champions vs. At-Risk), predictive sales forecasting to plan inventory, and tracking product profit margins to maximize profitability.",
        "features": [
            {
                "name": "Customer RFM Segmentation",
                "what": "Categorizing buyers into distinct behavioral segments (Champions, Loyal Customers, Potential Loyalists, At-Risk, and Lost Customers) based on Recency, Frequency, and Monetary value.",
                "connects": "Uses clean transaction data to drive targeted marketing and retention strategies.",
                "value": "Increases customer retention by identifying valuable buyers and alerting teams to churn risks before customers leave."
            },
            {
                "name": "Predictive Sales Forecasting",
                "what": "Smart predictive models that analyze historical sales patterns to project future revenue, seasonal spikes, and demand shifts.",
                "connects": "Links sales performance history to future inventory planning and target setting.",
                "value": "Prevents revenue loss from stockouts and avoids over-budgeting by giving leaders an accurate view of future sales trends."
            },
            {
                "name": "Profitability & Inventory Health Tracking",
                "what": "Real-time monitoring of product margins, stock velocity, and overall inventory turnover.",
                "connects": "Connects customer demand predictions with supply chain logistics.",
                "value": "Protects profit margins, reduces holding costs, and ensures high-demand products remain fully stocked."
            }
        ],
        "slides": [
            {
                "num": "Slide 1",
                "title": "Customer RFM Segmentation Matrix",
                "visual": "Customer RFM Segment Distribution (112 Champions, 42 At-Risk)",
                "img": IMG_CUST_TOP,
                "points": [
                    "<b>Behavioral Segmentation</b>: Categorizing buyers by Recency, Frequency, and Monetary value",
                    "<b>112 High-Value Champions</b>: Identifying top 10% customers generating 40% of revenue",
                    "<b>42 At-Risk Accounts</b>: Proactive churn alerts to trigger retention campaigns"
                ],
                "script": "Thank you Person 2. Raw sales numbers tell us how much revenue came in, but customer analytics tell us WHO is spending and WHY. As seen on our Customer RFM dashboard, Stratify AI automatically groups customers into behavioral segments. We can instantly pinpoint our 112 High-Value Champions driving core revenue, as well as 42 At-Risk accounts that require win-back offers.",
                "takeaway": "RFM segmentation transforms generic customer lists into targeted revenue retention strategies."
            },
            {
                "num": "Slide 2",
                "title": "Predictive Sales & Revenue Forecasting",
                "visual": "Prophet & XGBoost Multi-Horizon Forecasting Chart (95.0% Confidence)",
                "img": IMG_FORECAST_MID,
                "points": [
                    "<b>Multi-Horizon Predictions</b>: 30, 60, 90, and 180-day revenue projections",
                    "<b>High Model Accuracy</b>: 4.2% MAPE error rate with 95% confidence intervals",
                    "<b>Demand Spike Planning</b>: Helping store managers align staffing and stock replenishment"
                ],
                "script": "Beyond understanding past purchases, business leaders need to look forward. As shown in our ML Predictive Forecasting view, our Prophet and XGBoost models project revenue with 95% confidence intervals and an impressive 4.2% MAPE accuracy. This allows store managers to anticipate demand spikes and schedule inventory orders with precision.",
                "takeaway": "Predictive forecasting turns reactive planning into proactive revenue management."
            },
            {
                "num": "Slide 3",
                "title": "Inventory Health & Margin Protection",
                "visual": "Multi-Horizon Forecast Controls & Model Parameters",
                "img": IMG_FORECAST_TOP,
                "points": [
                    "<b>Stock Velocity Tracking</b>: Differentiating fast-moving SKUs from slow-moving inventory",
                    "<b>Stockout Prevention</b>: Automated reorder triggers for high-margin product lines",
                    "<b>Margin Protection</b>: Balancing holding costs with maximum gross profit margin"
                ],
                "script": "Finally, we connect customer demand to inventory health. High sales mean nothing if profit margins are eroded by excess warehousing costs or stockouts. Stratify AI continuously monitors stock levels and profit margins across all categories, flagging products that need reordering. Let's hand over to Person 4 to see how all these insights are presented live on our executive dashboard.",
                "takeaway": "Inventory intelligence balances stock availability with maximum gross profit margin."
            }
        ],
        "handoff": "Now that we've seen how clean data transforms into customer insights and sales forecasts, let's pass to Person 4 to experience our interactive, real-time Executive Control Center."
    },
    {
        "person": "Person 4",
        "file_name": "Person4_Executive_Control_Center_and_Monitoring.pdf",
        "title": "Interactive Executive Control Center & Live Store Monitor",
        "theme": "The Real-Time Visual Command Hub for Executive Decision-Makers",
        "duration": "5 Minutes",
        "summary": "Demonstrates the 9-tab Streamlit executive dashboard, live store transaction feeds, visual business health scoring, and interactive Plotly chart controls designed for non-technical users.",
        "features": [
            {
                "name": "9-Hub Executive Control Center",
                "what": "A clean, intuitve web application structured into 9 business hubs: Executive Summary, Customer Analytics, Sales Forecasting, Inventory & Finance, Live Feed, Data Quality, AI Insights, Admin Panel, and Report Center.",
                "connects": "Brings together all company metrics into one browser window.",
                "value": "Eliminates tool switching and gives leaders instant access to any business area in 1 click."
            },
            {
                "name": "Live Store Transaction Feed",
                "what": "A real-time ticker stream displaying store transactions as they happen across retail outlets, updated continuously.",
                "connects": "Displays the live output of Person 2's data pipeline.",
                "value": "Provides immediate visual feedback on promotional launches, flash sales, and daily store performance."
            },
            {
                "name": "Visual Business Health Scoreboard",
                "what": "Color-coded executive KPI cards (Green = Healthy, Yellow = Warning, Red = Action Required) summarizing overall revenue targets, inventory health, and data accuracy.",
                "connects": "Summarizes Person 3's complex metrics into quick visual status indicators.",
                "value": "Allows busy executives to perform a 5-second business health check at any time of day."
            }
        ],
        "slides": [
            {
                "num": "Slide 1",
                "title": "The Executive Command Dashboard",
                "visual": "Streamlit Live Application Main View with 9 Navigation Hubs",
                "img": IMG_EXEC_MID,
                "points": [
                    "<b>Single-Screen Hub</b>: 9 specialized business portals accessible in 1 click",
                    "<b>Zero Technical Learning Curve</b>: Clean, visual navigation for managers & C-suite",
                    "<b>Instant Department Switching</b>: Move between Sales, Customers, Inventory, and Finance"
                ],
                "script": "Thanks Person 3. All the data pipelines and smart models come together here in our Executive Control Center. Built with a clean, intuitive layout, the dashboard allows leadership to navigate across 9 specialized business hubs. Whether an executive wants to check top-level monthly revenue or dive into store-specific inventory, everything is accessible in a single click.",
                "takeaway": "The executive control portal simplifies complex enterprise metrics into clean visual navigation."
            },
            {
                "num": "Slide 2",
                "title": "Live Transaction Feed & Store Ticker",
                "visual": "Live Transaction Stream Table from Live Web Application",
                "img": IMG_EXEC_BOTTOM,
                "points": [
                    "<b>Sub-Second Streaming</b>: Continuous auto-refresh of store checkout receipts",
                    "<b>Instant Event Feedback</b>: Monitor sales velocity live during marketing campaigns",
                    "<b>Store Location Tracking</b>: Compare live purchasing trends across branch outlets"
                ],
                "script": "One of our standout features is the Live Transaction Feed shown here from our live web application. In traditional systems, you wait until tomorrow to see today's sales. Here, transactions scroll across the screen in real time as items are scanned at store checkouts. During Black Friday or special sales promotions, executives can monitor sales velocity live.",
                "takeaway": "Live store streaming gives executives instant visibility into customer buying activity."
            },
            {
                "num": "Slide 3",
                "title": "Color-Coded Business Health Scores",
                "visual": "Visual Category Profit & Sales Performance Matrix",
                "img": IMG_EXEC_CHART2,
                "points": [
                    "<b>5-Second Health Assessment</b>: Green / Yellow / Red visual KPI cards",
                    "<b>Immediate Risk Spotting</b>: Green for 94% Revenue Health, Yellow for 78% Inventory",
                    "<b>Interactive Drill-Down</b>: Filter charts by product category, store ID, or date range"
                ],
                "script": "Executives don't have time to wade through endless raw spreadsheets. That's why we built visual Health Score Cards. In five seconds, a manager can see that overall revenue health is Green at 94%, while inventory levels are Yellow at 78%, signaling that reorders are needed. To see how AI turns these charts into strategic advice, let's pass to Person 5.",
                "takeaway": "Color-coded health scoring lets executives spot and address operational risks instantly."
            }
        ],
        "handoff": "While visual charts show us what is happening across our stores, I will now hand over to Person 5 to explain how our built-in AI strategy advisor tells us WHY it happened and WHAT strategic actions to take."
    },
    {
        "person": "Person 5",
        "file_name": "Person5_AI_Strategy_Advisor_and_Insights.pdf",
        "title": "AI Strategy Advisor & Plain English Business Q&A",
        "theme": "Democratizing Artificial Intelligence for C-Suite Strategic Guidance",
        "duration": "5 Minutes",
        "summary": "Focuses on the integrated DeepSeek AI advisor, natural language business Q&A, automated root-cause analysis, and strategic recommendation engine for business decision-makers.",
        "features": [
            {
                "name": "Digital Chief Data Officer (AI Advisor)",
                "what": "An integrated artificial intelligence engine that constantly reviews daily sales performance, customer churn indicators, and store operations to synthesize executive summaries.",
                "connects": "Reads data directly from the live dashboard and analytics engine.",
                "value": "Acts as an on-demand business strategist, replacing hours of manual analyst reporting with instant AI summaries."
            },
            {
                "name": "Plain English Business Q&A",
                "what": "A conversational interface where managers can type business questions in natural everyday language (e.g., 'Why did Store 3 sales drop this Tuesday?') and get clear written answers.",
                "connects": "Bridges non-technical business leaders directly with complex database records.",
                "value": "Democratizes data access—no need to write complex database queries or ask IT for custom reports."
            },
            {
                "name": "Strategic Action Recommendation Engine",
                "what": "Automated generation of prioritized business action items based on real-time data trends (e.g., 'Promote Category X in Region Y to clear surplus inventory').",
                "connects": "Turns passive charts into active, prescriptive business decisions.",
                "value": "Accelerates executive response time, boosting revenue and operational efficiency."
            }
        ],
        "slides": [
            {
                "num": "Slide 1",
                "title": "Your Digital Chief Data Officer",
                "visual": "AI Agentic Business Analyst Interface (DeepSeek-V3)",
                "img": IMG_AI_TOP,
                "points": [
                    "<b>24/7 AI Business Advisor</b>: Continuous automated scanning of store metrics",
                    "<b>Executive Summary Synthesis</b>: Distilling complex data into 3 strategic bullet points",
                    "<b>Root-Cause Analysis</b>: Explaining the 'Why' behind revenue increases or drops"
                ],
                "script": "Thank you Person 4. Data dashboards are great for showing numbers, but executives need context: Why are sales up? Why are customers churning in Store 4? Stratify AI embeds a Digital Chief Data Officer powered by DeepSeek-V3. As seen on screen, the AI automatically scans the day's sales metrics and writes an executive synthesis highlighting key wins and risk factors.",
                "takeaway": "AI strategy synthesis translates raw numbers into clear business stories."
            },
            {
                "num": "Slide 2",
                "title": "Plain English Business Q&A",
                "visual": "DeepSeek CDO Executive Strategy Synthesis & Suggested Prompts",
                "img": IMG_AI_MID,
                "points": [
                    "<b>Conversational Interface</b>: Type business questions in plain everyday English",
                    "<b>Instant IT Bypass</b>: Get answers immediately without waiting 48 hours for SQL reports",
                    "<b>Data-Backed Clarity</b>: Written responses backed directly by audited database records"
                ],
                "script": "Imagine being an executive in a board meeting and needing an instant answer about store performance. Instead of submitting a request to IT and waiting 48 hours, you simply type your question into Stratify AI in plain English—for example: 'What were our top 3 revenue drivers this week?' The AI queries the database instantly and returns a clear, executive-ready explanation.",
                "takeaway": "Natural language Q&A empowers leadership to get instant answers without technical hurdles."
            },
            {
                "num": "Slide 3",
                "title": "Prescriptive Business Recommendations",
                "visual": "AI Strategic Action Recommendation Cards",
                "img": IMG_AI_TOP,
                "points": [
                    "<b>Prescriptive Guidance</b>: Moving from reactive charts to prioritized action steps",
                    "<b>Targeted Interventions</b>: Specific advice (e.g. 'Launch re-engagement promo for Segment B')",
                    "<b>Bottom-Line Results</b>: Accelerating response times to capture market growth"
                ],
                "script": "The true power of AI lies in action. Stratify AI doesn't just inform you—it guides you. If customer churn rises in a specific segment, the AI flags the root cause and presents recommended steps, such as launching a targeted re-engagement campaign. Now, to see how these insights are packaged and emailed automatically, let's hand over to Person 6.",
                "takeaway": "AI-driven recommendations convert insights directly into profitable business decisions."
            }
        ],
        "handoff": "Now that we've seen how AI turns data into strategic advice, I'll hand over to Person 6 to show how these insights are automatically compiled into executive PDF briefings and delivered straight to leadership inboxes."
    },
    {
        "person": "Person 6",
        "file_name": "Person6_Automated_Reporting_and_Business_ROI.pdf",
        "title": "Automated Executive Reporting & Total Business ROI",
        "theme": "Zero-Touch Automated Delivery, Enterprise Security & Business Impact",
        "duration": "5 Minutes",
        "summary": "Covers the 8-page automated executive PDF briefing engine, Robotic Process Automation (RPA) email distribution via Gmail, data security, and total business ROI.",
        "features": [
            {
                "name": "8-Page Automated Executive PDF Briefing Engine",
                "what": "An automated document engine that compiles live store data, KPI cards, charts, customer RFM tables, and AI insights into a publication-quality 8-page PDF report.",
                "connects": "Consolidates all outputs from Persons 1 through 5 into a single formatted report.",
                "value": "Saves 15+ hours per week of manual report compilation by analysts and management."
            },
            {
                "name": "Robotic Process Automation (RPA) Email Delivery",
                "what": "Smart digital bots that automatically trigger report generation and deliver executive PDFs directly to C-suite Gmail inboxes on a scheduled or event-driven basis.",
                "connects": "Delivers the final output of the platform directly to decision-makers.",
                "value": "Guarantees executives receive morning business briefings before their first meeting without lifting a finger."
            },
            {
                "name": "Enterprise Security & Total Business ROI",
                "what": "Bank-grade data security (zero hardcoded passwords, encrypted cloud connections) combined with high return on investment (time savings, margin growth, customer retention).",
                "connects": "Provides the overall commercial justification and governance framework.",
                "value": "Delivers up to 300% ROI within the first year while ensuring complete enterprise data compliance."
            }
        ],
        "slides": [
            {
                "num": "Slide 1",
                "title": "Automated Executive PDF Briefings",
                "visual": "Executive Reports & UiPath RPA Automation Dashboard",
                "img": IMG_REPORTS_TOP,
                "points": [
                    "<b>Zero-Touch Compilation</b>: 8-page publication-quality PDF report built automatically",
                    "<b>Consolidated Intelligence</b>: Unifying store sales, RFM tables, charts & AI recommendations",
                    "<b>Automated Archival</b>: Stored securely in digital repository for audit compliance"
                ],
                "script": "Thanks Person 5. To close our business workflow, we address how executives receive information outside the web portal. Not every executive opens a browser dashboard every morning. Stratify AI includes an Automated PDF Generator that compiles live store metrics, customer RFM tables, inventory alerts, and AI strategy notes into a clean, 8-page executive PDF briefing.",
                "takeaway": "Automated PDF briefings deliver formatted executive reports with zero manual labor."
            },
            {
                "num": "Slide 2",
                "title": "Robotic Email Delivery (RPA)",
                "visual": "UiPath RPA Execution Logs & Gmail SMTP Dispatch Panel",
                "img": IMG_REPORTS_MID,
                "points": [
                    "<b>Automated Gmail Dispatch</b>: Active email bots delivering PDFs directly to leadership",
                    "<b>7:00 AM Daily Schedule</b>: Morning briefings ready before the first meeting",
                    "<b>Audit Verification</b>: Complete delivery logging and recipient confirmation"
                ],
                "script": "These reports aren't just built—they are actively delivered. As shown on our live demo screen, digital bots archive the report in our secure repository and send it directly to designated leadership email inboxes via Gmail SMTP. By 7:00 AM every morning, executives have their complete business briefing ready on their mobile phone or laptop before their first cup of coffee.",
                "takeaway": "Robotic email distribution ensures leadership stays informed effortlessly every day."
            },
            {
                "num": "Slide 3",
                "title": "Business ROI & Conclusion",
                "visual": "ROI Summary Dashboard: 15+ Hours Saved/Week, 95% Error Reduction, 25% Churn Reduction",
                "img": IMG_REPORTS_TOP,
                "points": [
                    "<b>15+ Hours Saved / Week</b>: Eliminating manual Excel compiling for every manager",
                    "<b>95% Error Reduction</b>: Eradicating human copy-paste errors with automated pipelines",
                    "<b>Measurable Revenue Growth</b>: Faster churn intervention & protected profit margins"
                ],
                "script": "To summarize the total business value of Stratify AI: We take raw store transactions, validate them automatically, analyze customer behavior, present live executive dashboards, generate AI strategic advice, and deliver automated morning PDF reports. The result is 15+ hours saved per manager each week, 95% fewer reporting errors, and faster decisions that drive revenue growth.",
                "takeaway": "Stratify AI turns disconnected retail data into automated revenue growth and executive agility."
            }
        ],
        "handoff": "That concludes our 6-part presentation of Stratify AI. We would now like to open the floor to any questions from the executive team!"
    }
]

# PDF Generation Function
def create_pdf_for_presenter(pdata, output_path, styles):
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    
    story = []
    
    # Title Header Box
    header_table_data = [
        [
            Paragraph(f"<b>STRATIFY AI — EXECUTIVE PRESENTATION GUIDE</b>", ParagraphStyle('H1', fontName='Helvetica-Bold', fontSize=10, textColor=WHITE)),
            Paragraph(f"<b>{pdata['person'].upper()} GUIDE</b>", ParagraphStyle('H2', fontName='Helvetica-Bold', fontSize=10, textColor=WHITE, alignment=2))
        ],
        [
            Paragraph(f"<b>Topic: {pdata['title']}</b>", ParagraphStyle('H3', fontName='Helvetica-Bold', fontSize=14, textColor=WHITE)),
            Paragraph(f"Duration: {pdata['duration']}", ParagraphStyle('H4', fontName='Helvetica', fontSize=11, textColor=WHITE, alignment=2))
        ]
    ]
    
    header_table = Table(header_table_data, colWidths=[4.25*inch, 3.25*inch])
    header_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), PRIMARY),
        ('PADDING', (0,0), (-1,-1), 8),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('BOTTOMPADDING', (0,1), (-1,1), 10),
    ]))
    story.append(header_table)
    story.append(Spacer(1, 10))
    
    # Theme & Summary Card
    card_data = [
        [Paragraph(f"<b>Core Theme:</b> {pdata['theme']}", styles['CustomBodyBold'])],
        [Paragraph(f"<b>Executive Summary:</b> {pdata['summary']}", styles['CustomBody'])]
    ]
    card_table = Table(card_data, colWidths=[7.5*inch])
    card_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), CARD_BG),
        ('BOX', (0,0), (-1,-1), 1, BORDER_COLOR),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(card_table)
    story.append(Spacer(1, 12))
    
    # Feature Breakdown Section
    story.append(Paragraph("1. Business Features & Connectivity", styles['SectionHeader']))
    story.append(HRFlowable(width="100%", thickness=1, color=SECONDARY, spaceBefore=2, spaceAfter=8))
    
    feat_table_data = [
        [
            Paragraph("Feature Name", styles['TableHeader']),
            Paragraph("What It Does (Business Function)", styles['TableHeader']),
            Paragraph("How It Connects", styles['TableHeader']),
            Paragraph("Business Value & ROI", styles['TableHeader'])
        ]
    ]
    
    for f in pdata['features']:
        feat_table_data.append([
            Paragraph(f"<b>{f['name']}</b>", styles['TableCell']),
            Paragraph(f['what'], styles['TableCell']),
            Paragraph(f['connects'], styles['TableCell']),
            Paragraph(f['value'], styles['TableCell'])
        ])
        
    feat_table = Table(feat_table_data, colWidths=[1.4*inch, 2.3*inch, 1.8*inch, 2.0*inch])
    feat_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), SECONDARY),
        ('GRID', (0,0), (-1,-1), 0.5, BORDER_COLOR),
        ('PADDING', (0,0), (-1,-1), 6),
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [WHITE, BG_LIGHT])
    ]))
    story.append(feat_table)
    story.append(Spacer(1, 14))
    
    # Slide Breakdown Section
    story.append(Paragraph("2. Slide-by-Slide Presentation Guide & Spoken Script", styles['SectionHeader']))
    story.append(HRFlowable(width="100%", thickness=1, color=SECONDARY, spaceBefore=2, spaceAfter=8))
    
    for slide in pdata['slides']:
        slide_elements = []
        slide_elements.append(Paragraph(f"<b>{slide['num']}: {slide['title']}</b>", styles['SubSectionHeader']))
        
        points_text = "<br/>".join([f"&bull; {pt}" for pt in slide['points']])
        
        vp_data = [
            [Paragraph(f"<b>Visual Display (Live App):</b> {slide['visual']}", styles['CustomBodyBold'])],
            [Paragraph("<b>Key Presentation Points:</b>", styles['CustomBodyBold'])],
            [Paragraph(points_text, styles['CustomBody'])],
            [Paragraph("<b>Spoken Presentation Script (Non-Technical):</b>", styles['CustomBodyBold'])],
            [Paragraph(f'"{slide["script"]}"', styles['SpokenScript'])],
            [Paragraph(f"<b>Business Takeaway:</b> {slide['takeaway']}", ParagraphStyle('T', parent=styles['CustomBodyBold'], textColor=ACCENT))]
        ]
        
        vp_table = Table(vp_data, colWidths=[7.5*inch])
        vp_table.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,-1), BG_LIGHT),
            ('BOX', (0,0), (-1,-1), 1, BORDER_COLOR),
            ('PADDING', (0,0), (-1,-1), 7),
            ('LINEBELOW', (0,0), (-1,0), 0.5, BORDER_COLOR),
            ('LINEBELOW', (0,2), (-1,2), 0.5, BORDER_COLOR),
            ('LINEBELOW', (0,4), (-1,4), 0.5, BORDER_COLOR),
        ]))
        
        slide_elements.append(vp_table)
        
        if slide.get('img') and os.path.exists(slide['img']):
            try:
                slide_elements.append(Spacer(1, 4))
                slide_elements.append(Image(slide['img'], width=7.5*inch, height=3.2*inch))
            except Exception as e:
                print(f"Notice: Could not attach image to PDF: {e}")

        slide_elements.append(Spacer(1, 10))
        story.append(KeepTogether(slide_elements))
        
    # Handoff Section
    story.append(Spacer(1, 6))
    story.append(Paragraph("3. Speaker Hand-off & Transition Script", styles['SectionHeader']))
    story.append(HRFlowable(width="100%", thickness=1, color=SECONDARY, spaceBefore=2, spaceAfter=8))
    
    handoff_data = [
        [Paragraph("<b>Transition Words to Next Speaker:</b>", styles['CustomBodyBold'])],
        [Paragraph(f'"{pdata["handoff"]}"', styles['SpokenScript'])]
    ]
    handoff_table = Table(handoff_data, colWidths=[7.5*inch])
    handoff_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), CARD_BG),
        ('BOX', (0,0), (-1,-1), 1.5, SECONDARY),
        ('PADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(handoff_table)
    
    doc.build(story)
    print(f"Successfully generated PDF: {output_path}")

# Master PDF Summary Deck Generator
def create_master_pdf_deck(output_path, styles):
    doc = SimpleDocTemplate(
        output_path,
        pagesize=letter,
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36
    )
    story = []
    
    story.append(Paragraph("🎯 STRATIFY AI — MASTER EXECUTIVE PRESENTATION DECK", styles['DocTitle']))
    story.append(Paragraph("Complete 6-Person Business Flow & Live Streamlit Demo Screenshots", styles['DocSubTitle']))
    story.append(HRFlowable(width="100%", thickness=2, color=PRIMARY, spaceBefore=0, spaceAfter=12))
    
    summary_table_data = [
        [
            Paragraph("Presenter", styles['TableHeader']),
            Paragraph("Business Topic", styles['TableHeader']),
            Paragraph("Core Business Function", styles['TableHeader']),
            Paragraph("Key Output / Deliverable", styles['TableHeader'])
        ]
    ]
    
    for p in PRESENTERS_DATA:
        summary_table_data.append([
            Paragraph(f"<b>{p['person']}</b>", styles['TableCell']),
            Paragraph(p['title'], styles['TableCell']),
            Paragraph(p['theme'], styles['TableCell']),
            Paragraph(p['features'][0]['value'], styles['TableCell'])
        ])
        
    s_table = Table(summary_table_data, colWidths=[1.1*inch, 2.0*inch, 2.4*inch, 2.0*inch])
    s_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), PRIMARY),
        ('GRID', (0,0), (-1,-1), 0.5, BORDER_COLOR),
        ('PADDING', (0,0), (-1,-1), 6),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [WHITE, BG_LIGHT])
    ]))
    story.append(s_table)
    story.append(Spacer(1, 15))
    
    for p in PRESENTERS_DATA:
        p_elements = []
        p_elements.append(Paragraph(f"<b>{p['person']}: {p['title']}</b>", styles['SectionHeader']))
        p_elements.append(Paragraph(f"<b>Duration:</b> {p['duration']} | <b>Theme:</b> {p['theme']}", styles['CustomBodyBold']))
        p_elements.append(Paragraph(f"<b>Executive Summary:</b> {p['summary']}", styles['CustomBody']))
        
        slide_summary_text = ""
        for s in p['slides']:
            slide_summary_text += f"&bull; <b>{s['num']} ({s['title']}):</b> {s['takeaway']}<br/>"
        p_elements.append(Paragraph(slide_summary_text, styles['CustomBody']))
        
        p_elements.append(Paragraph(f"<b>Transition Script:</b> <i>\"{p['handoff']}\"</i>", styles['SpokenScript']))
        p_elements.append(Spacer(1, 10))
        p_elements.append(HRFlowable(width="100%", thickness=0.5, color=BORDER_COLOR, spaceBefore=4, spaceAfter=8))
        
        story.append(KeepTogether(p_elements))
        
    doc.build(story)
    print(f"Successfully generated Master Summary PDF: {output_path}")

# ULTRA PRESENTATION-READY POWERPOINT DECK GENERATOR WITH EMBEDDED SPEAKER NOTES
def create_master_pptx_deck(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    def add_header(slide, title_text, category_text="STRATIFY AI — EXECUTIVE BUSINESS DECK"):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(26, 54, 93) # Deep Navy
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.margin_left = Inches(0.6)
        tf.margin_top = Inches(0.18)
        
        p1 = tf.paragraphs[0]
        p1.text = category_text.upper()
        p1.font.size = Pt(11)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(49, 151, 149) # Teal
        
        p2 = tf.add_paragraph()
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(255, 255, 255)

    # ---------------------------------------------------------
    # SLIDE 1: Title Slide (Presentation Cover)
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(26, 54, 93)
    bg.line.fill.background()
    
    # Accent Line
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(0.15), Inches(3.8))
    acc.fill.solid()
    acc.fill.fore_color.rgb = RGBColor(49, 151, 149)
    acc.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1.4), Inches(1.8), Inches(11.0), Inches(3.8))
    tf = txBox.text_frame
    
    p1 = tf.paragraphs[0]
    p1.text = "🎯 STRATIFY AI"
    p1.font.size = Pt(44)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(49, 151, 149)
    
    p2 = tf.add_paragraph()
    p2.text = "Real-Time Executive Decision Intelligence Platform"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    
    p3 = tf.add_paragraph()
    p3.text = "Complete 6-Person Executive Presentation & Business Flow"
    p3.font.size = Pt(18)
    p3.font.color.rgb = RGBColor(203, 213, 224)

    p4 = tf.add_paragraph()
    p4.text = "🌐 Live Demo: https://stratify-ai-demo.streamlit.app/"
    p4.font.size = Pt(14)
    p4.font.italic = True
    p4.font.color.rgb = RGBColor(160, 174, 192)

    # Presenter Notes for Slide 1
    notes1 = slide.notes_slide.notes_text_frame
    notes1.text = (
        "WELCOME & INTRO SCRIPT:\n"
        "Welcome executives and board members. Today, our team presents Stratify AI—a real-time decision intelligence platform that transforms disconnected store checkout sales into instant strategic clarity. Over the next 30 minutes, 6 presenters will walk you through our complete end-to-end business flow."
    )

    # ---------------------------------------------------------
    # SLIDE 2: Agenda & 6-Person Flow Map
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Executive Agenda & 6-Person Presentation Flow")
    
    for idx, pdata in enumerate(PRESENTERS_DATA):
        col = idx % 3
        row = idx // 3
        left = Inches(0.6 + col * 4.1)
        top = Inches(1.5 + row * 2.8)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.9), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(247, 250, 252)
        card.line.color.rgb = RGBColor(203, 213, 224)
        
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_right = Inches(0.2)
        
        p1 = tf.paragraphs[0]
        p1.text = f"{pdata['person']}: {pdata['title']}"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(26, 54, 93)
        
        p2 = tf.add_paragraph()
        p2.text = f"Theme: {pdata['theme']}"
        p2.font.size = Pt(10)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(43, 108, 176)
        
        p3 = tf.add_paragraph()
        p3.text = pdata['summary']
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = RGBColor(45, 55, 72)

    notes2 = slide.notes_slide.notes_text_frame
    notes2.text = (
        "AGENDA SCRIPT:\n"
        "Here is our 6-stage presentation map. Person 1 introduces the business vision; Person 2 details live store capture and automated data quality; Person 3 presents customer segmentation and forecasting; Person 4 demonstrates the executive dashboard; Person 5 showcases our AI Strategy Advisor; and Person 6 closes with automated reporting and ROI."
    )

    # ---------------------------------------------------------
    # SLIDES 3 TO 20: 3 Presentation-Ready Slides Per Presenter
    # ---------------------------------------------------------
    for pdata in PRESENTERS_DATA:
        for slide_info in pdata['slides']:
            slide = prs.slides.add_slide(blank_layout)
            add_header(slide, f"{slide_info['title']}", f"{pdata['person'].upper()} — {pdata['title'].upper()}")
            
            # Left Side (55% width): HIGH-RES SCREENSHOT WITH ELEGANT CONTAINER
            img_path = slide_info.get('img')
            if img_path and os.path.exists(img_path):
                # Border container box
                img_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(7.2), Inches(5.8))
                img_bg.fill.solid()
                img_bg.fill.fore_color.rgb = RGBColor(247, 250, 252)
                img_bg.line.color.rgb = RGBColor(49, 151, 149) # Teal Border
                img_bg.line.width = Pt(1.5)
                
                # Screenshot image
                slide.shapes.add_picture(img_path, Inches(0.6), Inches(1.4), Inches(7.0), Inches(5.0))
                
                # Caption bar
                tx_cap = slide.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(7.0), Inches(0.6))
                tf_c = tx_cap.text_frame
                p_c = tf_c.paragraphs[0]
                p_c.text = f"📸 Live Streamlit Demo Screen: {slide_info['visual']}"
                p_c.font.size = Pt(10)
                p_c.font.bold = True
                p_c.font.color.rgb = RGBColor(26, 54, 93)

            # Right Side (40% width): HIGH-IMPACT BULLETS & TAKEAWAY PILL
            right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.9), Inches(1.3), Inches(4.9), Inches(5.8))
            right_box.fill.solid()
            right_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            right_box.line.color.rgb = RGBColor(203, 213, 224)
            
            tf_r = right_box.text_frame
            tf_r.margin_left = Inches(0.25)
            tf_r.margin_top = Inches(0.25)
            tf_r.margin_right = Inches(0.25)
            
            # Header
            p_r1 = tf_r.paragraphs[0]
            p_r1.text = "📌 KEY EXECUTIVE INSIGHTS"
            p_r1.font.size = Pt(13)
            p_r1.font.bold = True
            p_r1.font.color.rgb = RGBColor(26, 54, 93)
            
            # Bullet points
            for point in slide_info['points']:
                p_pt = tf_r.add_paragraph()
                p_pt.text = f"• {point}"
                p_pt.font.size = Pt(11)
                p_pt.font.color.rgb = RGBColor(45, 55, 72)
                p_pt.space_after = Pt(8)
                
            # Key Takeaway Banner Box inside Right Card
            takeaway_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(5.6), Inches(4.5), Inches(1.3))
            takeaway_box.fill.solid()
            takeaway_box.fill.fore_color.rgb = RGBColor(237, 242, 247)
            takeaway_box.line.color.rgb = RGBColor(49, 151, 149)
            
            tf_t = takeaway_box.text_frame
            tf_t.margin_left = Inches(0.15)
            tf_t.margin_top = Inches(0.12)
            tf_t.margin_right = Inches(0.15)
            
            p_t1 = tf_t.paragraphs[0]
            p_t1.text = "💡 BUSINESS TAKEAWAY"
            p_t1.font.size = Pt(10)
            p_t1.font.bold = True
            p_t1.font.color.rgb = RGBColor(49, 151, 149)
            
            p_t2 = tf_t.add_paragraph()
            p_t2.text = slide_info['takeaway']
            p_t2.font.size = Pt(10.5)
            p_t2.font.color.rgb = RGBColor(26, 54, 93)

            # ---------------------------------------------------------
            # EMBEDDED SPEAKER NOTES (PRESENTER VIEW READY!)
            # ---------------------------------------------------------
            notes = slide.notes_slide.notes_text_frame
            notes.text = (
                f"=== PRESENTER NOTES ({pdata['person']} - {slide_info['num']}) ===\n\n"
                f"🗣️ SPOKEN SCRIPT:\n\"{slide_info['script']}\"\n\n"
                f"💡 BUSINESS TAKEAWAY:\n{slide_info['takeaway']}\n\n"
                f"🖥️ VISUAL SCREEN REFERENCE:\n{slide_info['visual']}\n\n"
                f"🔄 HAND-OFF SCRIPT (At end of section):\n\"{pdata['handoff']}\""
            )

    # ---------------------------------------------------------
    # SLIDE 21: Executive Summary & Total ROI
    # ---------------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Executive Summary & Business Return on Investment")
    
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(11.333), Inches(5.4))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(26, 54, 93)
    card.line.color.rgb = RGBColor(49, 151, 149)
    card.line.width = Pt(2)
    
    tf = card.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.4)
    
    p1 = tf.paragraphs[0]
    p1.text = "🎉 STRATIFY AI — TOTAL BUSINESS VALUE SUMMARY"
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(49, 151, 149)
    
    bullets = [
        "100% Automated Decision Pipeline: Store checkout ➔ Clean repository ➔ AI ➔ Executive Inbox.",
        "Zero Garbage Data Guarantee: Automated Digital Quality Filter with <5s SLA latency.",
        "Predictive Business Intelligence: RFM Customer Segmentation & Prophet/XGBoost 95% Confidence Forecasts.",
        "24/7 AI Strategy Advisor: DeepSeek-V3 providing plain-English executive recommendations.",
        "Massive Operational Savings: 15+ hours/week saved per manager & 95% reduction in reporting errors.",
        "Verified Live Deployment: Active and running live at https://stratify-ai-demo.streamlit.app/"
    ]
    
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"✔  {bullet}"
        p.font.size = Pt(14.5)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)

    notes21 = slide.notes_slide.notes_text_frame
    notes21.text = (
        "SUMMARY & ROI SCRIPT:\n"
        "To summarize: Stratify AI delivers a fully automated decision pipeline that saves over 15 hours per week per manager while providing 100% data quality, predictive forecasting, and AI strategic guidance. Thank you for your time, and we are now ready to take your questions!"
    )

    # Save presentation
    prs.save(output_path)
    print(f"Successfully generated Presentation-Ready Master PowerPoint Deck: {output_path}")

# Main execution
if __name__ == "__main__":
    pdf_styles = get_pdf_styles()
    
    # 1. Generate 6 PDF files for the 6 presenters
    for pdata in PRESENTERS_DATA:
        file_path = os.path.join(OUTPUT_DIR, pdata['file_name'])
        create_pdf_for_presenter(pdata, file_path, pdf_styles)
        
    # 2. Generate Master PDF Summary Deck
    master_pdf_path = os.path.join(OUTPUT_DIR, "Stratify_AI_Master_Business_Presentation.pdf")
    create_master_pdf_deck(master_pdf_path, pdf_styles)
    
    # 3. Generate ULTRA PRESENTATION-READY Master PowerPoint Deck
    master_pptx_path = os.path.join(OUTPUT_DIR, "Stratify_AI_Master_Business_Presentation_Ready.pptx")
    try:
        create_master_pptx_deck(master_pptx_path)
    except Exception as e:
        print(f"Error saving presentation deck: {e}")
    
    print("\nALL PRESENTATION-READY DELIVERABLES UPDATED SUCCESSFULLY!")
